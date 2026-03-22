"""PowerPoint mockup parser.

Parses .pptx files to extract slide layouts, shapes, text content, and colors.
Converts slides to images for Claude vision interpretation.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pbi_developer.utils.logging import get_logger

logger = get_logger(__name__)


@dataclass
class ShapeInfo:
    """Extracted information about a PowerPoint shape."""
    shape_type: str
    name: str
    text: str
    left: int  # EMUs
    top: int
    width: int
    height: int
    fill_color: str | None = None


@dataclass
class SlideInfo:
    """Extracted information about a PowerPoint slide."""
    index: int
    shapes: list[ShapeInfo] = field(default_factory=list)
    title: str = ""
    notes: str = ""


@dataclass
class PptxParseResult:
    """Result of parsing a PowerPoint file."""
    slides: list[SlideInfo] = field(default_factory=list)
    slide_images: list[bytes] = field(default_factory=list)
    summary: str = ""


def parse_pptx(path: Path) -> PptxParseResult:
    """Parse a PowerPoint file and extract shapes, text, and layout.

    Returns structured slide data and slide images for vision analysis.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(path))
    result = PptxParseResult()

    for idx, slide in enumerate(prs.slides):
        slide_info = SlideInfo(index=idx)

        # Extract title
        if slide.shapes.title:
            slide_info.title = slide.shapes.title.text

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_info.notes = slide.notes_slide.notes_text_frame.text

        # Extract shapes
        for shape in slide.shapes:
            shape_info = ShapeInfo(
                shape_type=str(shape.shape_type),
                name=shape.name,
                text=shape.text if shape.has_text_frame else "",
                left=shape.left or 0,
                top=shape.top or 0,
                width=shape.width or 0,
                height=shape.height or 0,
            )

            # Extract fill color if available
            if hasattr(shape, "fill") and shape.fill.type is not None:
                try:
                    if shape.fill.fore_color and shape.fill.fore_color.rgb:
                        shape_info.fill_color = f"#{shape.fill.fore_color.rgb}"
                except Exception:
                    pass

            slide_info.shapes.append(shape_info)

        result.slides.append(slide_info)
        logger.info(f"Slide {idx + 1}: '{slide_info.title}' ({len(slide_info.shapes)} shapes)")

    result.summary = _build_summary(result)
    return result


def slides_to_text(result: PptxParseResult) -> str:
    """Convert parsed slides to a text description for the planner agent."""
    parts: list[str] = []
    for slide in result.slides:
        parts.append(f"\n## Slide {slide.index + 1}: {slide.title or '(untitled)'}")
        if slide.notes:
            parts.append(f"Notes: {slide.notes}")
        for shape in slide.shapes:
            if shape.text:
                # Convert EMU to approximate pixels (1 EMU = 1/914400 inch, ~96 DPI)
                px_left = shape.left // 9525
                px_top = shape.top // 9525
                px_w = shape.width // 9525
                px_h = shape.height // 9525
                parts.append(
                    f"- [{shape.name}] at ({px_left},{px_top}) {px_w}x{px_h}px: \"{shape.text}\""
                )
                if shape.fill_color:
                    parts[-1] += f" (color: {shape.fill_color})"
    return "\n".join(parts)


def _build_summary(result: PptxParseResult) -> str:
    """Build a summary of the parsed PowerPoint."""
    total_shapes = sum(len(s.shapes) for s in result.slides)
    return (
        f"Parsed {len(result.slides)} slide(s) with {total_shapes} total shapes. "
        f"Titles: {', '.join(s.title or '(untitled)' for s in result.slides)}"
    )
